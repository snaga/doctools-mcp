from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import win32com.client

def count_slides(input_path: str) -> dict:
    """
    Get the total number of slides in a PowerPoint file.
    
    Args:
        input_path: Path to the PPTX file.
        
    Returns:
        dict: Result containing slide count.
    """
    try:
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")
            
        prs = Presentation(input_path)
        return {
            "status": "success",
            "slide_count": len(prs.slides)
        }
        
    except Exception as e:
        return {"status": "error", "detail": str(e)}

def extract_text_as_markdown(input_path: str, output_path: str = None, start_slide: int = 1, end_slide: int = None) -> dict:
    """
    Extract text from a PowerPoint file and save it as a Markdown file.
    
    Args:
        input_path: Path to the PPTX file.
        output_path: Path to the output Markdown file (optional).
        start_slide: Starting slide number (1-based, default 1).
        end_slide: Ending slide number (1-based, default None=last slide).
        
    Returns:
        dict: Result containing the absolute path to the generated Markdown file.
    """
    try:
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")
            
        prs = Presentation(input_path)
        total_slides = len(prs.slides)
        
        if end_slide is None:
            end_slide = total_slides
            
        # Validation
        if start_slide < 1:
            raise ValueError(f"Start slide must be >= 1. Got: {start_slide}")
        if end_slide > total_slides:
            raise ValueError(f"End slide {end_slide} exceeds total slides {total_slides}.")
        if start_slide > end_slide:
            raise ValueError(f"Start slide {start_slide} cannot be greater than end slide {end_slide}.")
            
        extracted_text = []
        for i in range(start_slide - 1, end_slide):
            slide = prs.slides[i]
            slide_text = []
            
            # Title
            title_shape = getattr(slide.shapes, "title", None)
            if title_shape and getattr(title_shape, "text", None):
                slide_text.append(f"### {title_shape.text}")
                
            # Content
            for shape in slide.shapes:
                # Avoid duplicating title
                if title_shape == shape:
                    continue

                # Text Frame (Text Box, Placeholder)
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text.append(f"- {paragraph.text}")
                
                # Table
                elif getattr(shape, "has_table", False):
                    table = shape.table
                    # Header
                    headers = []
                    # Simple heuristic: first row is header
                    first_row = True
                    for row in table.rows:
                        row_cells = [cell.text_frame.text.replace("\n", " ").strip() for cell in row.cells]
                        row_str = "| " + " | ".join(row_cells) + " |"
                        slide_text.append(row_str)
                        if first_row:
                            sep_str = "| " + " | ".join(["---"] * len(row_cells)) + " |"
                            slide_text.append(sep_str)
                            first_row = False
                    slide_text.append("") # Empty line after table

                # Image / Picture
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    name = getattr(shape, "name", "Picture")
                    slide_text.append(f"\n[Image: {name}]")

                # Chart
                elif getattr(shape, "has_chart", False):
                    title = ""
                    try:
                        title = shape.chart.chart_title.text_frame.text
                    except:
                        title = "Chart"
                    slide_text.append(f"\n[Chart: {title}]")
            
            # Notes
            if getattr(slide, "has_notes_slide", False):
                notes_slide = getattr(slide, "notes_slide", None)
                if notes_slide:
                    notes_tf = getattr(notes_slide, "notes_text_frame", None)
                    if notes_tf and getattr(notes_tf, "text", "").strip():
                         slide_text.append(f"\n> Notes: {notes_tf.text}")

            content = "\n".join(slide_text)
            extracted_text.append(f"## Slide {i+1}\n{content}")
        
        if output_path is None:
            output_path = os.path.splitext(input_path)[0] + ".md"
            
        markdown_content = "\n\n".join(extracted_text)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(markdown_content)
                
        return {
            "status": "success",
            "output_path": os.path.abspath(output_path)
        }
    except Exception as e:
        return {"status": "error", "detail": str(e)}

def extract_slides(input_path: str, output_path: str, start_slide: int, end_slide: int) -> dict:
    """
    Extract specified slides from a PowerPoint file and save as a new file.
    
    Args:
        input_path: Path to the source PPTX file.
        output_path: Path to save the extracted slides.
        start_slide: Starting slide number (1-based).
        end_slide: Ending slide number (1-based).
        
    Returns:
        dict: Result of the operation including status and output path.
    """
    try:
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")
            
        # python-pptx doesn't have a simple "copy slide" API to a new presentation easily
        # without losing styles/masters. A common workaround is to delete unwanted slides
        # from a copy of the presentation.
        
        prs = Presentation(input_path)
        total_slides = len(prs.slides)
        
        # Validation
        if start_slide < 1:
            raise ValueError(f"Start slide must be >= 1. Got: {start_slide}")
        if end_slide > total_slides:
            raise ValueError(f"End slide {end_slide} exceeds total slides {total_slides}.")
        if start_slide > end_slide:
            raise ValueError(f"Start slide {start_slide} cannot be greater than end slide {end_slide}.")

        # Indices to keep (0-based)
        keep_indices = set(range(start_slide - 1, end_slide))
        
        # We delete slides in reverse order to avoid index shifting issues
        # XML extraction involves creating a new list of slide IDs
        
        # python-pptx deletes via xml manipulation under the hood when we know the rId
        # simpler approach: iterate all slides backwards, if index not in keep_indices, delete.
        
        # Accessing slides collection directly
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        
        for i in reversed(range(total_slides)):
            if i not in keep_indices:
                xml_slides.remove(slides[i])
                
        prs.save(output_path)
            
        return {
            "status": "success",
            "message": "PPTX slides extracted successfully.",
            "output_path": output_path
        }
        
    except Exception as e:
        return {"status": "error", "detail": str(e)}

def export_slides_to_images(
    input_path: str,
    output_dir: str = None,
    slides: list[int] = None,
    width: int = 1280,
    height: int = 720
) -> dict:
    """
    Export specified slides from a PowerPoint file as images (PNG).
    
    Args:
        input_path: Path to the PPTX file.
        output_dir: Directory to save the images (optional).
        slides: List of 1-based slide numbers to export (optional, default all).
        width: Width of the exported images (default 1280).
        height: Height of the exported images (default 720).
        
    Returns:
        dict: Result containing status and list of output image paths.
    """
    powerpoint = None
    presentation = None
    try:
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")
            
        abs_input_path = os.path.abspath(input_path)
        
        if output_dir is None:
            output_dir = os.path.splitext(abs_input_path)[0] + "_images"
        
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
        abs_output_dir = os.path.abspath(output_dir)
        
        # Initialize PowerPoint COM
        # Note: This requires Microsoft PowerPoint to be installed on Windows.
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        
        # Open presentation (WithWindow=False for background)
        # 1st arg: FileName, 2nd: ReadOnly, 3rd: Untitled, 4th: WithWindow
        presentation = powerpoint.Presentations.Open(abs_input_path, WithWindow=False)
        
        total_slides = presentation.Slides.Count
        
        if slides is None:
            slides = list(range(1, total_slides + 1))
            
        output_paths = []
        for slide_num in slides:
            if 1 <= slide_num <= total_slides:
                slide = presentation.Slides(slide_num)
                image_path = os.path.join(abs_output_dir, f"slide_{slide_num:03d}.png")
                # Export(Path, FilterName, ScaleWidth, ScaleHeight)
                slide.Export(image_path, "PNG", width, height)
                output_paths.append(image_path)
            else:
                # Skip invalid slide numbers
                continue
                
        return {
            "status": "success",
            "output_paths": output_paths
        }
        
    except Exception as e:
        return {"status": "error", "detail": str(e)}
    finally:
        if presentation:
            try:
                presentation.Close()
            except:
                pass
        if powerpoint:
            try:
                powerpoint.Quit()
            except:
                pass
